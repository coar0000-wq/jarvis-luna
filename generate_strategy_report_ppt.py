#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""JARVIS 5일 주기 전략분석 PowerPoint 생성 및 이메일 발송 스크립트."""

import json
import os
import smtplib
from datetime import datetime, timedelta, timezone
from email.encoders import encode_base64
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path
from typing import Any
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


class StrategyReportGenerator:
    """최근 5일 데이터를 바탕으로 전략 보고서 PPTX를 생성하고 발송한다."""

    OUTPUT_DIR = Path("data")
    RECIPIENT_EMAIL = "coar0000@naver.com"

    NAVY = (15, 35, 66)
    BLUE = (37, 99, 235)
    SKY = (219, 234, 254)
    WHITE = (255, 255, 255)
    TEXT = (31, 41, 55)
    MUTED = (75, 85, 99)
    GREEN = (5, 150, 105)

    def __init__(self) -> None:
        self.now = datetime.now(timezone.utc)
        self.output_dir = self.OUTPUT_DIR
        self.recipient_email = os.getenv("RECIPIENT_EMAIL", self.RECIPIENT_EMAIL)
        self.sender_email = os.getenv("SENDER_EMAIL", "").strip()
        self.sender_password = os.getenv("EMAIL_PASSWORD", "")
        self.email_error: str | None = None

    @staticmethod
    def _load_json(path: Path, default: dict[str, Any]) -> dict[str, Any]:
        """JSON 파일을 읽고, 파일 부재 또는 형식 오류 시 기본값을 사용한다."""
        try:
            with path.open("r", encoding="utf-8-sig") as file:
                loaded = json.load(file)
                return loaded if isinstance(loaded, dict) else default
        except (FileNotFoundError, OSError, json.JSONDecodeError) as error:
            print(f"⚠️ 데이터 파일을 읽지 못해 기본값을 사용합니다: {path} ({error})")
            return default

    @staticmethod
    def _number(value: Any, default: float = 0.0) -> float:
        """통화·백분율 문자열을 계산 가능한 숫자로 변환한다."""
        try:
            return float(str(value).replace(",", "").replace("%", "").strip())
        except (TypeError, ValueError):
            return default

    def fetch_shopify_order_metrics(self) -> dict[str, Any] | None:
        """최근 5일의 Shopify 주문을 읽어 매출·이행 지표로 집계한다."""
        shop_domain = os.getenv("SHOPIFY_SHOP_URL", "").strip().removeprefix("https://").removeprefix("http://").rstrip("/")
        access_token = os.getenv("SHOPIFY_ACCESS_TOKEN", "").strip()
        api_version = os.getenv("SHOPIFY_API_VERSION", "2026-07").strip()
        if not shop_domain or not access_token:
            print("ℹ️ Shopify Secrets가 없어 기존 운영 요약 JSON을 전략 분석 입력으로 사용합니다.")
            return None

        since = (self.now - timedelta(days=5)).strftime("%Y-%m-%dT%H:%M:%SZ")
        query = """
        query RecentOrders($query: String!, $after: String) {
          orders(first: 250, after: $after, query: $query, sortKey: PROCESSED_AT, reverse: true) {
            nodes {
              processedAt
              cancelledAt
              displayFinancialStatus
              displayFulfillmentStatus
              currentTotalPriceSet { shopMoney { amount currencyCode } }
            }
            pageInfo { hasNextPage endCursor }
          }
        }
        """
        orders: list[dict[str, Any]] = []
        after: str | None = None
        try:
            while True:
                payload = json.dumps(
                    {"query": query, "variables": {"query": f"processed_at:>={since} status:any", "after": after}}
                ).encode("utf-8")
                request = Request(
                    f"https://{shop_domain}/admin/api/{api_version}/graphql.json",
                    data=payload,
                    headers={
                        "Content-Type": "application/json",
                        "X-Shopify-Access-Token": access_token,
                    },
                    method="POST",
                )
                with urlopen(request, timeout=30) as response:
                    response_data = json.loads(response.read().decode("utf-8"))
                if response_data.get("errors"):
                    raise RuntimeError(response_data["errors"])
                result = response_data.get("data", {}).get("orders", {})
                orders.extend(result.get("nodes", []))
                page_info = result.get("pageInfo", {})
                if not page_info.get("hasNextPage"):
                    break
                after = page_info.get("endCursor")
        except (HTTPError, URLError, OSError, ValueError, RuntimeError) as error:
            print(f"⚠️ Shopify 주문 API를 읽지 못해 기존 운영 요약 JSON을 사용합니다: {error}")
            return None

        active_orders = [order for order in orders if not order.get("cancelledAt")]
        fulfilled_states = {"FULFILLED", "SHIPPED"}
        pending_states = {"UNFULFILLED", "PARTIAL", "ON_HOLD", "SCHEDULED"}
        fulfilled_orders = [order for order in active_orders if order.get("displayFulfillmentStatus") in fulfilled_states]
        pending_orders = [order for order in active_orders if order.get("displayFulfillmentStatus") in pending_states]
        total_revenue = sum(
            self._number(order.get("currentTotalPriceSet", {}).get("shopMoney", {}).get("amount"))
            for order in active_orders
        )
        currency = next(
            (
                order.get("currentTotalPriceSet", {}).get("shopMoney", {}).get("currencyCode")
                for order in active_orders
                if order.get("currentTotalPriceSet", {}).get("shopMoney", {}).get("currencyCode")
            ),
            "USD",
        )
        today = self.now.date().isoformat()
        shipped_today = sum(
            1
            for order in fulfilled_orders
            if str(order.get("processedAt", "")).startswith(today)
        )
        fulfillment_rate = (len(fulfilled_orders) / len(active_orders) * 100) if active_orders else 0.0
        return {
            "source": "shopify_admin_graphql",
            "collected_at": self.now.isoformat(),
            "period_start": since,
            "orders_5d": len(active_orders),
            "revenue_5d": round(total_revenue, 2),
            "currency": currency,
            "average_order_value": round(total_revenue / len(active_orders), 2) if active_orders else 0.0,
            "fulfilled_orders_5d": len(fulfilled_orders),
            "pending_orders": len(pending_orders),
            "cancelled_orders_5d": len(orders) - len(active_orders),
            "shipped_today": shipped_today,
            "fulfillment_rate": round(fulfillment_rate, 1),
        }

    def _write_operational_snapshots(self, metrics: dict[str, Any], shopify_snapshot: dict[str, Any]) -> None:
        """실시간 Shopify 집계를 기존 대시보드 JSON에도 반영한다."""
        shopify_snapshot.update(
            {
                "store_status": shopify_snapshot.get("store_status", "온라인"),
                "sync_status": f"Shopify Admin API 동기화: {self.now.strftime('%Y-%m-%d %H:%M UTC')}",
                "orders_5d": metrics["orders_5d"],
                "revenue_5d": metrics["revenue_5d"],
                "currency": metrics["currency"],
                "average_order_value": metrics["average_order_value"],
            }
        )
        order_snapshot = {
            "pending_orders": metrics["pending_orders"],
            "shipped_today": metrics["shipped_today"],
            "fulfillment_rate": f"{metrics['fulfillment_rate']:.1f}%",
            "status": "정상 운영" if metrics["fulfillment_rate"] >= 95 else "배송 점검 필요",
            "source": "Shopify Admin API",
            "updated_at": self.now.isoformat(),
        }
        with (self.output_dir / "shopify_team.json").open("w", encoding="utf-8") as file:
            json.dump(shopify_snapshot, file, ensure_ascii=False, indent=2)
        with (self.output_dir / "order_team.json").open("w", encoding="utf-8") as file:
            json.dump(order_snapshot, file, ensure_ascii=False, indent=2)

    def collect_5day_data(self) -> dict[str, Any]:
        """상품 데이터와 Shopify·주문 운영 지표를 최근 5일 분석 입력으로 수집한다."""
        daiso_data = self._load_json(
            self.output_dir / "daiso_products.json",
            {"total_count": 0, "products": []},
        )
        dropshipping_data = self._load_json(
            self.output_dir / "global_daiso_dropshipping.json",
            {"revenue_forecast": {}},
        )
        shopify_snapshot = self._load_json(
            self.output_dir / "shopify_team.json",
            {"visitors_today": 0, "conversion_rate": "0%"},
        )
        order_snapshot = self._load_json(
            self.output_dir / "order_team.json",
            {"pending_orders": 0, "shipped_today": 0, "fulfillment_rate": "0%"},
        )
        live_metrics = self.fetch_shopify_order_metrics()
        if live_metrics:
            self._write_operational_snapshots(live_metrics, shopify_snapshot)
            shopify_data = {**shopify_snapshot, **live_metrics}
            data_quality = "Shopify Admin API 및 저장소 수집 데이터 기준"
        else:
            shopify_data = {
                "source": "dashboard_snapshot",
                "visitors_today": shopify_snapshot.get("visitors_today", 0),
                "conversion_rate": shopify_snapshot.get("conversion_rate", "0%"),
                "orders_5d": shopify_snapshot.get("orders_5d"),
                "revenue_5d": shopify_snapshot.get("revenue_5d"),
                "currency": shopify_snapshot.get("currency", "USD"),
                "average_order_value": shopify_snapshot.get("average_order_value"),
                "pending_orders": order_snapshot.get("pending_orders", 0),
                "shipped_today": order_snapshot.get("shipped_today", 0),
                "fulfillment_rate": self._number(order_snapshot.get("fulfillment_rate")),
            }
            data_quality = "저장소의 Shopify·주문 운영 스냅샷과 수집 데이터 기준"
        return {
            "timestamp": self.now.isoformat(),
            "collection_period": "5 days",
            "daiso": daiso_data,
            "dropshipping": dropshipping_data,
            "shopify": shopify_data,
            "data_quality": data_quality,
        }

    def analyze_strategy(self, data: dict[str, Any]) -> dict[str, Any]:
        """Shopify 주문·이행 지표와 상품 데이터를 함께 반영해 전략 우선순위를 구성한다."""
        total_products = data.get("daiso", {}).get("total_count", 0)
        shopify = data.get("shopify", {})
        conversion_rate = self._number(shopify.get("conversion_rate"))
        fulfillment_rate = self._number(shopify.get("fulfillment_rate"))
        pending_orders = int(self._number(shopify.get("pending_orders")))
        orders_5d = shopify.get("orders_5d")
        revenue_5d = shopify.get("revenue_5d")
        currency = shopify.get("currency", "USD")
        recommendations: list[dict[str, str]] = []

        if pending_orders > 0:
            recommendations.append(
                {
                    "title": "주문·배송 병목 우선 해소",
                    "description": f"현재 미처리 주문 {pending_orders}건을 배송 단계·재고·고객 안내 기준으로 분류하고, 처리 기한을 먼저 확정합니다.",
                    "impact": "배송 지연·취소 위험을 낮추고 고객 경험을 보호",
                }
            )
        if conversion_rate and conversion_rate < 2.5:
            recommendations.append(
                {
                    "title": "전환율 개선 실험",
                    "description": f"현재 전환율 {conversion_rate:.1f}%를 기준으로 상품 상세 페이지, 배송 메시지, 가격·번들 제안을 우선 A/B 점검합니다.",
                    "impact": "동일 유입 대비 주문 전환 개선",
                }
            )
        if fulfillment_rate and fulfillment_rate < 95:
            recommendations.append(
                {
                    "title": "풀필먼트 품질 회복",
                    "description": f"현재 주문 이행률 {fulfillment_rate:.1f}%의 하락 요인을 재고·배송사·처리 리드타임으로 분해해 개선합니다.",
                    "impact": "배송 완결률 및 재구매 신뢰도 개선",
                }
            )
        if revenue_5d is not None:
            recommendations.append(
                {
                    "title": "최근 주문 매출 기반 객단가 최적화",
                    "description": f"최근 5일 Shopify 매출 {currency} {self._number(revenue_5d):,.2f}와 주문 수를 기준으로 번들·업셀·무료배송 임계값을 조정합니다.",
                    "impact": "객단가와 매출 효율 개선",
                }
            )
        recommendations.extend(
            [
                {
                    "title": "고마진 카테고리 집중",
                    "description": "수익성이 높은 핵심 상품군의 재고·광고 효율을 집중 점검하고, 주문 성과가 확인된 품목의 노출을 확대합니다.",
                    "impact": "운영 자원을 수익 기회가 큰 상품에 집중",
                },
                {
                    "title": "글로벌 시장 확대",
                    "description": "주문과 전환 성과가 안정적인 상품부터 국가별 상세 페이지와 캠페인 메시지를 현지화합니다.",
                    "impact": "검증된 상품 중심의 시장 도달 범위 확대",
                },
                {
                    "title": "가격·재고 운영 최적화",
                    "description": "최근 주문·배송 데이터와 경쟁 가격을 함께 검토해 품절 위험과 과도한 할인 없이 목표 수익성을 유지합니다.",
                    "impact": "수익성·재고 회전·배송 안정성의 균형 개선",
                },
            ]
        )
        revenue_label = (
            f"{currency} {self._number(revenue_5d):,.2f}"
            if revenue_5d is not None
            else "Shopify 연동 대기"
        )
        orders_label = f"{int(self._number(orders_5d))}건" if orders_5d is not None else "Shopify 연동 대기"
        return {
            "timestamp": self.now.isoformat(),
            "period": "최근 5일",
            "key_metrics": {
                "total_products": total_products,
                "shopify_orders_5d": orders_label,
                "shopify_revenue_5d": revenue_label,
                "fulfillment_rate": f"{fulfillment_rate:.1f}%" if fulfillment_rate else "운영 데이터 대기",
                "conversion_rate": f"{conversion_rate:.1f}%" if conversion_rate else "운영 데이터 대기",
                "pending_orders": f"{pending_orders}건",
            },
            "shopify_operational": shopify,
            "strategy_recommendations": recommendations[:5],
            "data_quality": data["data_quality"],
        }

    @staticmethod
    def _set_fill(shape: Any, color: tuple[int, int, int]) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()

    @staticmethod
    def _set_run_font(run: Any, size: float, color: tuple[int, int, int], bold: bool = False) -> None:
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)

    def _add_textbox(
        self,
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: float,
        color: tuple[int, int, int],
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    ) -> Any:
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = textbox.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.04)
        frame.margin_right = Inches(0.04)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        frame.vertical_anchor = vertical_anchor
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._set_run_font(run, font_size, color, bold)
        return textbox

    def _add_slide_header(self, slide: Any, title: str, page_number: int) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(*self.WHITE)

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.18),
        )
        self._set_fill(accent, self.BLUE)
        self._add_textbox(slide, title, 0.72, 0.45, 10.8, 0.45, 28, self.NAVY, bold=True)
        self._add_textbox(
            slide,
            f"JARVIS 전략 보고서  |  {page_number}",
            10.9,
            0.55,
            1.7,
            0.25,
            9,
            self.MUTED,
            align=PP_ALIGN.RIGHT,
        )

    def generate_ppt(self, analysis: dict[str, Any]) -> tuple[Path, Path]:
        """분석 결과를 실제 .pptx 및 메타데이터 JSON으로 저장한다."""
        self.output_dir.mkdir(parents=True, exist_ok=True)
        stamp = self.now.strftime("%Y%m%d_%H%M%S")
        ppt_path = self.output_dir / f"strategy_report_{stamp}.pptx"
        metadata_path = self.output_dir / f"strategy_report_{stamp}.json"

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        blank_layout = presentation.slide_layouts[6]

        # 1. 표지
        slide = presentation.slides.add_slide(blank_layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(*self.NAVY)
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(5.95), Inches(13.333), Inches(1.55),
        )
        self._set_fill(banner, self.BLUE)
        self._add_textbox(slide, "JARVIS", 0.82, 0.78, 3.0, 0.6, 22, self.SKY, bold=True)
        self._add_textbox(slide, "5일 주기 전략분석", 0.78, 1.55, 10.9, 0.9, 36, self.WHITE, bold=True)
        self._add_textbox(
            slide,
            "데이터 기반 성장 전략 및 실행 우선순위",
            0.82,
            2.62,
            9.8,
            0.45,
            19,
            self.SKY,
        )
        self._add_textbox(
            slide,
            self.now.astimezone().strftime("보고서 생성일  %Y-%m-%d  |  %H:%M %Z"),
            0.84,
            6.43,
            8.5,
            0.35,
            13,
            self.WHITE,
        )

        # 2. 핵심 지표
        slide = presentation.slides.add_slide(blank_layout)
        self._add_slide_header(slide, "핵심 지표", 2)
        metrics = analysis["key_metrics"]
        cards = [
            ("총 상품 수", f"{metrics['total_products']}개"),
            ("최근 5일 주문", metrics["shopify_orders_5d"]),
            ("최근 5일 Shopify 매출", metrics["shopify_revenue_5d"]),
            ("주문 이행률", metrics["fulfillment_rate"]),
        ]
        for index, (label, value) in enumerate(cards):
            col = index % 2
            row = index // 2
            left = 0.9 + col * 6.15
            top = 1.45 + row * 2.5
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top), Inches(5.35), Inches(1.85),
            )
            self._set_fill(card, self.SKY)
            self._add_textbox(slide, label, left + 0.34, top + 0.32, 4.6, 0.32, 15, self.MUTED, bold=True)
            self._add_textbox(slide, value, left + 0.34, top + 0.82, 4.6, 0.58, 25, self.NAVY, bold=True)
        self._add_textbox(
            slide,
            "지표는 최근 5일 Shopify 주문·배송 운영 데이터와 저장소 수집 데이터를 바탕으로 산출되었습니다.",
            0.95,
            6.55,
            11.4,
            0.3,
            10,
            self.MUTED,
        )

        # 3~7. 전략 제안
        for page_number, recommendation in enumerate(analysis["strategy_recommendations"], start=3):
            slide = presentation.slides.add_slide(blank_layout)
            self._add_slide_header(slide, f"전략 {page_number - 2}. {recommendation['title']}", page_number)

            strategy_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.9), Inches(1.45), Inches(11.55), Inches(2.1),
            )
            self._set_fill(strategy_box, self.SKY)
            self._add_textbox(slide, "실행 제안", 1.3, 1.82, 2.0, 0.3, 15, self.BLUE, bold=True)
            self._add_textbox(slide, recommendation["description"], 1.3, 2.25, 10.7, 0.92, 21, self.TEXT)

            impact_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.9), Inches(4.08), Inches(11.55), Inches(1.28),
            )
            self._set_fill(impact_box, (236, 253, 245))
            self._add_textbox(slide, "기대 효과", 1.3, 4.38, 1.6, 0.28, 15, self.GREEN, bold=True)
            self._add_textbox(slide, recommendation["impact"], 3.05, 4.32, 8.8, 0.38, 18, self.TEXT, bold=True)

            self._add_textbox(
                slide,
                "다음 5일 동안 담당자·예산·측정 지표를 지정해 우선순위에 따라 실행하십시오.",
                1.0,
                6.32,
                11.2,
                0.32,
                11,
                self.MUTED,
                align=PP_ALIGN.CENTER,
            )

        presentation.save(ppt_path)

        metadata = {
            "filename": str(ppt_path),
            "title": f"JARVIS 5일 주기 전략분석 {self.now.strftime('%Y-%m-%d')}",
            "created_at": self.now.isoformat(),
            "analysis": analysis,
        }
        with metadata_path.open("w", encoding="utf-8") as file:
            json.dump(metadata, file, ensure_ascii=False, indent=2)

        if not ppt_path.is_file() or ppt_path.stat().st_size == 0:
            raise RuntimeError("PPTX 파일 생성에 실패했습니다.")
        return ppt_path, metadata_path

    def send_email(self, ppt_path: Path, analysis: dict[str, Any]) -> bool | None:
        """PPTX를 첨부해 전송한다. 메일 자격 증명이 없으면 전송을 건너뛴다."""
        if not self.sender_email or not self.sender_password:
            print("⚠️ SENDER_EMAIL 또는 EMAIL_PASSWORD가 설정되지 않아 이메일 전송을 건너뜁니다.")
            return None

        if not ppt_path.is_file():
            raise FileNotFoundError(f"첨부할 PPTX 파일이 없습니다: {ppt_path}")

        message = MIMEMultipart()
        message["From"] = self.sender_email
        message["To"] = self.recipient_email
        message["Subject"] = f"JARVIS 5일 전략분석 보고서 | {self.now.strftime('%Y-%m-%d')}"

        metrics = analysis["key_metrics"]
        recommendations = analysis["strategy_recommendations"][:3]
        strategy_lines = "\n".join(
            f"{index}. {item['title']} — {item['description']}"
            for index, item in enumerate(recommendations, start=1)
        )
        body = f"""안녕하세요.

JARVIS가 최근 5일 데이터를 기반으로 전략분석 보고서를 생성했습니다.

[주요 지표]
- 총 상품 수: {metrics['total_products']}개
- 최근 5일 Shopify 주문: {metrics['shopify_orders_5d']}
- 최근 5일 Shopify 매출: {metrics['shopify_revenue_5d']}
- 주문 이행률: {metrics['fulfillment_rate']}
- 전환율: {metrics['conversion_rate']}
- 미처리 주문: {metrics['pending_orders']}

[핵심 전략]
{strategy_lines}

상세 내용은 첨부된 PowerPoint 보고서를 확인해 주세요.

JARVIS
{self.now.strftime('%Y-%m-%d %H:%M UTC')}
"""
        message.attach(MIMEText(body, "plain", "utf-8"))

        with ppt_path.open("rb") as attachment:
            part = MIMEBase(
                "application",
                "vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            part.set_payload(attachment.read())
        encode_base64(part)
        part.add_header("Content-Disposition", "attachment", filename=ppt_path.name)
        message.attach(part)

        try:
            with smtplib.SMTP_SSL("smtp.naver.com", 465, timeout=30) as server:
                server.login(self.sender_email, self.sender_password)
                server.send_message(message)
            print(f"✅ 이메일 전송 완료: {self.recipient_email}")
            return True
        except (smtplib.SMTPException, OSError) as error:
            self.email_error = str(error)
            print(f"❌ 이메일 전송 실패: {self.email_error}")
            return False

    def write_log(self, ppt_path: Path, metadata_path: Path, email_status: bool | None, status: str) -> None:
        """실행 결과를 JSON Lines 형식의 로그 파일에 추가한다."""
        log_path = self.output_dir / "strategy_report_log.jsonl"
        log_entry = {
            "timestamp": self.now.isoformat(),
            "status": status,
            "ppt_file": str(ppt_path),
            "metadata_file": str(metadata_path),
            "recipient": self.recipient_email,
            "email_sent": email_status is True,
            "email_delivery_skipped": email_status is None,
            "email_error": self.email_error,
        }
        with log_path.open("a", encoding="utf-8") as file:
            file.write(json.dumps(log_entry, ensure_ascii=False) + "\n")

    def write_pages_snapshot(
        self,
        analysis: dict[str, Any],
        ppt_path: Path,
        metadata_path: Path,
        email_status: bool | None,
        status: str,
    ) -> Path:
        """GitHub Pages 대시보드가 읽는 최신 전략 보고서 요약을 저장한다."""
        snapshot_path = self.output_dir / "latest_strategy_report.json"
        email_delivery = "sent" if email_status is True else "skipped" if email_status is None else "retry_needed"
        snapshot = {
            "schema_version": 1,
            "title": f"JARVIS 5일 전략분석 | {self.now.strftime('%Y-%m-%d')}",
            "generated_at": self.now.isoformat(),
            "period": analysis["period"],
            "status": status,
            "key_metrics": analysis["key_metrics"],
            "shopify_operational": analysis["shopify_operational"],
            "strategy_recommendations": analysis["strategy_recommendations"],
            "pptx_url": ppt_path.name,
            "metadata_url": metadata_path.name,
            "email_delivery": email_delivery,
        }
        with snapshot_path.open("w", encoding="utf-8") as file:
            json.dump(snapshot, file, ensure_ascii=False, indent=2)
        return snapshot_path

    def run(self) -> None:
        """전체 보고서 생성, 이메일 발송, 결과 기록을 수행한다."""
        print(f"\n🎯 JARVIS 5일 전략분석 시작 ({self.now.isoformat()})")
        print("=" * 60)

        print("📊 데이터 수집 중...")
        data = self.collect_5day_data()
        print(f"✅ {data['daiso'].get('total_count', 0)}개 상품 데이터 로드 완료")

        print("🔍 전략분석 수행 중...")
        analysis = self.analyze_strategy(data)
        print(f"✅ {len(analysis['strategy_recommendations'])}개 전략 추천안 생성 완료")

        print("📄 실제 PowerPoint 생성 중...")
        ppt_path, metadata_path = self.generate_ppt(analysis)
        print(f"✅ PPTX 생성 완료: {ppt_path} ({ppt_path.stat().st_size:,} bytes)")
        print(f"✅ 메타데이터 저장 완료: {metadata_path}")

        print("📧 이메일 발송 중...")
        email_status = self.send_email(ppt_path, analysis)

        if email_status is True:
            status = "completed"
        elif email_status is None:
            status = "completed_without_email"
        else:
            status = "completed_with_email_error"

        snapshot_path = self.write_pages_snapshot(
            analysis,
            ppt_path,
            metadata_path,
            email_status,
            status,
        )
        self.write_log(ppt_path, metadata_path, email_status, status)
        print(f"✅ Pages 대시보드 요약 저장 완료: {snapshot_path}")
        if email_status is False:
            print("⚠️ 보고서와 Pages 데이터는 생성됐지만 이메일 전송은 재시도가 필요합니다.")
        print("=" * 60)
        print("🎉 JARVIS 5일 전략분석 완료!\n")


if __name__ == "__main__":
    StrategyReportGenerator().run()
